"""
Universal RAG Document Indexer
Supports: .txt, .pdf, .docx, .xlsx, .pptx

Install dependencies:
    pip install chromadb sentence-transformers anthropic
    pip install pypdf python-docx openpyxl python-pptx
"""

import os
import json
import chromadb
from chromadb.utils import embedding_functions
from anthropic import Anthropic
from pathlib import Path

# ── Optional imports (graceful degradation if library missing) ─────────────────
"""
try:
    from pypdf import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("Warning: pypdf not installed. PDF support disabled. Run: pip install pypdf")
"""

try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("Warning: python-docx not installed. DOCX support disabled. Run: pip install python-docx")

try:
    import openpyxl
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    print("Warning: openpyxl not installed. XLSX support disabled. Run: pip install openpyxl")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("Warning: python-pptx not installed. PPTX support disabled. Run: pip install python-pptx")


# ══════════════════════════════════════════════════════════════════════════════
# PARSERS — each returns plain text extracted from the file
# ══════════════════════════════════════════════════════════════════════════════

def parse_txt(filepath: str) -> str:
    """Plain text — just read it."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_pdf(filepath: str) -> str:
    pass
    """
    PDF — extract text page by page.
    Each page is separated by a marker so chunks don't bleed across pages.
    
    if not PDF_SUPPORT:
        raise RuntimeError("pypdf not installed")

    reader = PdfReader(filepath)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i+1}]\n{text.strip()}")
    return "\n\n".join(pages)
    """

def parse_docx(filepath: str) -> str:
    """
    DOCX — extract paragraphs and tables.
    Heading styles are preserved as text markers.
    """
    if not DOCX_SUPPORT:
        raise RuntimeError("python-docx not installed")

    doc = DocxDocument(filepath)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Preserve heading context so chunks know their section
        if para.style and para.style.name.startswith("Heading"):
            parts.append(f"\n## {text}")
        else:
            parts.append(text)

    # Also extract tables — each cell becomes part of a row string
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(f"[Table row] {row_text}")

    return "\n".join(parts)

def parse_xlsx(filepath: str) -> str:
    """
    XLSX — extract each sheet as structured text.
    Row/column context is preserved so the LLM can reason about the data.
    """
    if not XLSX_SUPPORT:
        raise RuntimeError("openpyxl not installed")

    wb = openpyxl.load_workbook(filepath, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = None

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            # Skip fully empty rows
            if all(cell is None for cell in row):
                continue

            row_values = [str(cell) if cell is not None else "" for cell in row]

            # Treat first non-empty row as header
            if headers is None:
                headers = row_values
                rows.append("HEADERS: " + " | ".join(headers))
            else:
                # Pair header → value for better semantic retrieval
                paired = " | ".join(
                    f"{h}: {v}" for h, v in zip(headers, row_values) if v
                )
                if paired:
                    rows.append(f"ROW {i}: {paired}")

        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    return "\n\n".join(sheets)

def parse_pptx(filepath: str) -> str:
    """
    PPTX — extract text from each slide.
    Slide number and title are preserved for context.
    """
    if not PPTX_SUPPORT:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"[Slide {i+1}]"]

        # Get slide title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"Title: {slide.shapes.title.text.strip()}")

        # Extract all text frames on the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue  # already handled above
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_parts.append(text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Speaker notes] {notes_text}")

        if len(slide_parts) > 1:  # more than just the slide number
            slides.append("\n".join(slide_parts))

    return "\n\n".join(slides)


# ══════════════════════════════════════════════════════════════════════════════
# PARSER REGISTRY — maps file extension to parser function
# ══════════════════════════════════════════════════════════════════════════════

PARSERS = {
    ".txt":  parse_txt,
    ".md":   parse_txt,   # markdown is plain text
    ".pdf":  parse_pdf,   # requires pypdf - NOT WORKING YET
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".xls":  parse_xlsx,
    ".pptx": parse_pptx,
}


# ══════════════════════════════════════════════════════════════════════════════
# CHUNKER
# ══════════════════════════════════════════════════════════════════════════════

def chunk_text(
    text: str,
    source: str,
    doc_type: str,
    chunk_size: int = 150,
    overlap: int = 30
) -> tuple[list[str], list[str], list[dict]]:
    """
    Split text into overlapping word-based chunks.
    Returns parallel lists: (chunks, ids, metadatas)
    """
    words = text.split()
    chunks, ids, metas = [], [], []

    for i, start in enumerate(range(0, len(words), chunk_size - overlap)):
        chunk = " ".join(words[start:start + chunk_size])
        if chunk.strip():
            chunk_id = f"{Path(source).stem}_{doc_type}_chunk_{i}"
            chunks.append(chunk)
            ids.append(chunk_id)
            metas.append({
                "source":   source,
                "doc_type": doc_type,
                "chunk_idx": i,
            })

    return chunks, ids, metas


# ══════════════════════════════════════════════════════════════════════════════
# RAG INDEXER CLASS
# ══════════════════════════════════════════════════════════════════════════════

class UniversalRAGIndexer:
    def __init__(self, db_path: str = "./chroma_db", collection_name: str = "documents"):
        self.client = chromadb.PersistentClient(path=db_path)
        self.embedder = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2"   # small, fast, fully offline
        )
        self.collection = self.client.get_or_create_collection(
            name=collection_name,
            embedding_function=self.embedder
        )

        print(f"ChromaDB ready at '{db_path}' — collection: '{collection_name}'")
        print(f"Documents already indexed: {self.collection.count()}\n")

    # ── Index a single file ───────────────────────────────────────────────────
    def index_file(self, filepath: str) -> int:
        path = Path(filepath)
        ext  = path.suffix.lower()

        if ext not in PARSERS:
            print(f"  Skipped (unsupported type): {filepath}")
            return 0

        print(f"  Parsing  [{ext}] {path.name} ...")
        try:
            text = PARSERS[ext](filepath)
        except Exception as e:
            print(f"  ERROR parsing {filepath}: {e}")
            return 0

        if not text.strip():
            print(f"  WARNING: No text extracted from {filepath}")
            return 0

        chunks, ids, metas = chunk_text(text, filepath, ext.lstrip("."))

        # Upsert — safe to re-index the same file
        self.collection.upsert(documents=chunks, ids=ids, metadatas=metas)
        print(f"  Indexed  {len(chunks)} chunks from {path.name}")
        return len(chunks)

    # ── Index an entire folder ────────────────────────────────────────────────
    def index_folder(self, folder_path: str) -> int:
        total = 0
        folder = Path(folder_path)
        files  = [f for f in folder.rglob("*") if f.suffix.lower() in PARSERS]

        print(f"Found {len(files)} supported files in '{folder_path}'")
        for f in files:
            total += self.index_file(str(f))

        print(f"\nDone. Total chunks indexed: {total}")
        return total


    # ── Show collection stats ─────────────────────────────────────────────────
    def stats(self):
        count = self.collection.count()
        print(f"\nCollection stats:")
        print(f"  Total chunks: {count}")
        if count > 0:
            # Show breakdown by doc type
            results = self.collection.get(include=["metadatas"])
            from collections import Counter
            types = Counter(m["doc_type"] for m in results["metadatas"])
            for dtype, n in types.most_common():
                print(f"  .{dtype}: {n} chunks")
    
    # ── Get list of indexed documents ──────────────────────────────────────────
    def get_indexed_documents(self):
        results = self.collection.get(include=["metadatas"])
        return [m["source"] for m in results["metadatas"]]
    

# ══════════════════════════════════════════════════════════════════════════════
# DEMO
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    rag = UniversalRAGIndexer(db_path="./test_chroma_db", collection_name="docs_1")

    # --- Index individual files of different types ---
    print("=== INDEXING ===")
    #rag.index_file("firewall_policy.txt")           # plain text
    #rag.index_file("audit_report.pdf")             # PDF
    #rag.index_file("audit_report.docx")             # Word document
    #rag.index_file("vulnerability_tracker.xlsx")    # Excel spreadsheet
    #rag.index_file("extract-tracker.xlsx")          # Excel spreadsheet
    #rag.index_file("cis_audit_presentation.pptx")        # PowerPoint

    # --- Or index an entire folder at once ---
    # rag.index_folder("./documents/")

    #rag.stats()
    print("\n=== INDEXED DOCUMENTS ===")
    print(rag.get_indexed_documents())
