# ── Optional imports (graceful degradation if library missing) ─────────────────
"""Some parsers require extra libraries. We try to import them here and set flags.
The main parsing functions will check these flags and raise errors if the library is missing.
try:
    from pypdf import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("Warning: pypdf not installed. PDF support disabled. Run: pip install pypdf")
"""
try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("Warning: python-docx not installed. DOCX support disabled. Run: pip install python-docx")

try:
    import openpyxl
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    print("Warning: openpyxl not installed. XLSX support disabled. Run: pip install openpyxl")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("Warning: python-pptx not installed. PPTX support disabled. Run: pip install python-pptx")

# ══════════════════════════════════════════════════════════════════════════════
# PARSERS — each returns plain text extracted from the file
# ══════════════════════════════════════════════════════════════════════════════

def parse_txt(filepath: str) -> str:
    """Plain text — just read it."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_pdf(filepath: str) -> str:
    pass

    """
    PDF — extract text page by page.
    Each page is separated by a marker so chunks don't bleed across pages.

    if not PDF_SUPPORT:
        raise RuntimeError("pypdf not installed")

    reader = PdfReader(filepath)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i+1}]\n{text.strip()}")
    return "\n\n".join(pages)
    """


def parse_docx(filepath: str) -> str:
    """
    DOCX — extract paragraphs and tables.
    Heading styles are preserved as text markers.
    """
    if not DOCX_SUPPORT:
        raise RuntimeError("python-docx not installed")

    doc = DocxDocument(filepath)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Preserve heading context so chunks know their section
        if para.style.name.startswith("Heading"):
            parts.append(f"\n## {text}")
        else:
            parts.append(text)

    # Also extract tables — each cell becomes part of a row string
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(f"[Table row] {row_text}")

    return "\n".join(parts)


def parse_xlsx(filepath: str) -> str:
    """
    XLSX — extract each sheet as structured text.
    Row/column context is preserved so the LLM can reason about the data.
    """
    if not XLSX_SUPPORT:
        raise RuntimeError("openpyxl not installed")

    wb = openpyxl.load_workbook(filepath, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = None

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            # Skip fully empty rows
            if all(cell is None for cell in row):
                continue

            row_values = [str(cell) if cell is not None else "" for cell in row]

            # Treat first non-empty row as header
            if headers is None:
                headers = row_values
                rows.append("HEADERS: " + " | ".join(headers))
            else:
                # Pair header → value for better semantic retrieval
                paired = " | ".join(
                    f"{h}: {v}" for h, v in zip(headers, row_values) if v
                )
                if paired:
                    rows.append(f"ROW {i}: {paired}")

        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    return "\n\n".join(sheets)


def parse_pptx(filepath: str) -> str:
    """
    PPTX — extract text from each slide.
    Slide number and title are preserved for context.
    """
    if not PPTX_SUPPORT:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"[Slide {i+1}]"]

        # Get slide title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"Title: {slide.shapes.title.text.strip()}")

        # Extract all text frames on the slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue  # already handled above
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_parts.append(text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Speaker notes] {notes_text}")

        if len(slide_parts) > 1:  # more than just the slide number
            slides.append("\n".join(slide_parts))

    return "\n\n".join(slides)